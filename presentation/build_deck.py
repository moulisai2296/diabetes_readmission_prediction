"""Build the demo deck (PPTX) — decisions & outcomes, stage by stage.

Strict layout grid, bottle-green palette, a native shapes-and-arrows architecture
diagram (no raster image). Numbers come from final_model_report.json + the Stage-3
comparison. No API/endpoint detail — the live demo covers the product itself.

  uv run python presentation/build_deck.py   # -> presentation/readmission_demo.pptx
"""

from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
HERE = ROOT / "presentation"

# palette (matches the UI)
ACCENT = RGBColor(0x0B, 0x6B, 0x4F)
DEEP = RGBColor(0x07, 0x4D, 0x39)
BG = RGBColor(0xD7, 0xE9, 0xDD)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT = RGBColor(0x16, 0x26, 0x1D)
MUTED = RGBColor(0x5A, 0x6F, 0x63)
PANEL = RGBColor(0xEE, 0xF5, 0xF0)
CHIP = RGBColor(0xCF, 0xE6, 0xD8)
LINE = RGBColor(0xB6, 0xD2, 0xC1)

SW, SH = Inches(13.333), Inches(7.5)
R = json.loads((ROOT / "src/models/final_model_report.json").read_text())
TA, T, CM = R["test"]["at_threshold"], R["test"], R["cost_model"]


def deck() -> Presentation:
    p = Presentation()
    p.slide_width, p.slide_height = SW, SH
    return p


def _bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def rect(slide, l, t, w, h, fill, line=None, rounded=True, shadow=False):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(l), Inches(t), Inches(w), Inches(h))
    shp.shadow.inherit = False
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(1.25)
    return shp


def textbox(slide, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.text_frame.word_wrap = True
    tb.text_frame.vertical_anchor = anchor
    return tb.text_frame


def para(tf, text, size, color, bold=False, first=False, align=PP_ALIGN.LEFT,
         after=6, bullet=False, indent=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.text = ("•  " if bullet and not indent else "–  " if indent else "") + text
    p.font.size = Pt(size); p.font.color.rgb = color; p.font.bold = bold
    p.alignment = align; p.space_after = Pt(after)
    if indent:
        p.level = 1
    return p


def head(slide, title, n=None):
    bar = rect(slide, 0, 0, 13.333, 0.95, ACCENT, rounded=False)
    tf = bar.text_frame; tf.margin_left = Inches(0.55); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(tf, title, 25, WHITE, bold=True, first=True)


def title_slide(p):
    s = p.slides.add_slide(p.slide_layouts[6]); _bg(s, ACCENT)
    rect(s, 0.7, 2.25, 0.18, 2.4, CHIP)  # accent stripe
    tf = textbox(s, 1.1, 2.3, 11.3, 2.4)
    para(tf, "Diabetes 30-Day Readmission Prediction", 40, WHITE, bold=True, first=True)
    para(tf, "Decisions, outcomes, and the system — stage by stage", 20, CHIP, after=0)
    tf2 = textbox(s, 1.1, 6.5, 11.3, 0.6)
    para(tf2, "Care-team decision support   ·   UCI Diabetes 130-US Hospitals (1999–2008)",
         13, CHIP, first=True)
    return s


def stage_slide(p, title, decisions, out_title, outcomes):
    """Consistent template: decisions (left) → outcome card (right)."""
    s = p.slides.add_slide(p.slide_layouts[6]); _bg(s, BG)
    head(s, title)
    tf = textbox(s, 0.6, 1.3, 7.0, 5.7)
    para(tf, "KEY DECISIONS", 14, ACCENT, bold=True, first=True, after=10)
    for d in decisions:
        para(tf, d.strip(), 15 if not d.startswith("  ") else 13,
             TEXT if not d.startswith("  ") else MUTED,
             bullet=not d.startswith("  "), indent=d.startswith("  "), after=7)
    rect(s, 8.0, 1.35, 4.73, 5.5, WHITE, line=LINE)
    rect(s, 8.0, 1.35, 4.73, 0.55, ACCENT)
    cap = rect(s, 8.0, 1.35, 4.73, 0.55, ACCENT)
    cap.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(cap.text_frame, "  " + out_title, 14, WHITE, bold=True, first=True)
    tf2 = textbox(s, 8.25, 2.1, 4.25, 4.6)
    for i, o in enumerate(outcomes):
        para(tf2, o.strip(), 14, TEXT, bullet=True, first=(i == 0), after=9)
    return s


def build():
    p = deck()
    title_slide(p)

    # 2. Problem & framing decisions
    s = p.slides.add_slide(p.slide_layouts[6]); _bg(s, BG)
    head(s, "The Problem & Framing Decisions")
    tf = textbox(s, 0.6, 1.3, 12.1, 5.6)
    para(tf, "Goal: at discharge, flag each diabetic patient likely to be readmitted within "
         "30 days, so care teams can target follow-up.", 18, TEXT, first=True, after=16)
    para(tf, "Target = binary “readmitted < 30 days” (≈11% of patients) — matches the real "
         "action (flag for follow-up).", 16, TEXT, bullet=True, after=10)
    para(tf, "Accuracy is the wrong metric — predicting “never readmit” scores 89%. We judge "
         "on recall, precision, calibration, and net savings.", 16, TEXT, bullet=True, after=10)
    para(tf, "The model is ~20% of the work; the system around it — reproducible pipeline, "
         "serving, monitoring, governance — is the other 80%.", 16, TEXT, bullet=True, after=10)
    para(tf, "Output is decision support, never an automatic action: a clinician sees the risk "
         "and the reasons, and decides.", 16, TEXT, bullet=True)

    # 3. Architecture (native shapes)
    architecture_slide(p)

    # 4. Stage 1 — Data
    stage_slide(p, "Stage 1 — Data Engineering",
        ["Missing values coded as “?” → treated as NaN, not blanks",
         "Drop weight (~97% missing); justify payer_code / specialty",
         "A missing A1C/glucose test is information →",
         "  kept as “NotMeasured”, never imputed to normal",
         "Exclude encounters that can’t be readmitted",
         "  (expired / hospice discharge)",
         "Define the target once, in one place"],
        "Outcome",
        ["99,340 clean encounters (~70k patients)",
         "Validated cleaning pipeline (pandera), DVC-tracked",
         "Reproducible: one command rebuilds the data",
         "Missingness preserved as signal, not noise"])

    # 5. Stage 2 — Features
    stage_slide(p, "Stage 2 — Feature Engineering",
        ["Group ICD-9 diag_1/2/3 into clinical buckets",
         "  (Circulatory, Respiratory, Diabetes, …)",
         "Roll up utilization: total_prior_visits, active meds,",
         "  medication changes; ordinal age",
         "Collapse rare specialty / payer categories to “Other”",
         "ONE shared transform used at train and serve time"],
        "Outcome",
        ["48 model features from ~26 raw inputs",
         "Server derives engineered features — caller sends facts",
         "No training/serving skew by construction",
         "Every feature + rationale logged in FEATURES.md"])

    # 6. Stage 3 — Modeling (table)
    modeling_slide(p)

    # 7. Stage 4 — Packaging
    stage_slide(p, "Stage 4 — Packaging & Deployment",
        ["Export model + feature spec as one matched artifact",
         "  (serving never touches MLflow or training data)",
         "Containerize the service; whole stack in one command",
         "Bundle metrics + dashboards with the service",
         "Cloud Run as the deploy target (scales to zero)",
         "Keep the previous image for one-command rollback"],
        "Outcome",
        ["Self-contained image — model travels with the code",
         "One command runs API + monitoring locally",
         "Same image deploys unchanged to the cloud",
         "Simple, documented rollback path"])

    # 8. Stage 5 — Monitoring
    stage_slide(p, "Stage 5 — Observability & Monitoring",
        ["Log every prediction (inputs, score, version, time)",
         "Track service health: latency, errors, request rate",
         "Detect drift vs the training distribution",
         "  (data drift + prediction-score drift)",
         "Alert a human on errors / unusual flag rate",
         "Define a concrete retraining trigger"],
        "Outcome",
        ["Live drift report (recent traffic vs training)",
         "Dashboards for service & model behaviour",
         "Alerts fire when something moves",
         "Retrain when: ≥30% feature drift, recall < 0.55",
         "on new labels, or quarterly"])

    # 9. Stage 6 — Governance
    stage_slide(p, "Stage 6 — Governance & Re-evaluation",
        ["Audit fairness across age, gender, race",
         "  (per-group recall & selection rate, not averages)",
         "Explain globally and per patient (SHAP)",
         "Trace lineage: data + code + model version",
         "Keep the human in the loop for borderline cases",
         "Write it down: model card + reflection"],
        "Outcome",
        ["Gender essentially fair (recall gap 0.045)",
         "Race / age gaps disclosed (small subgroups, young)",
         "Top drivers: prior inpatient visits, discharge type",
         "Model card + fairness report generated from code"])

    # 10. Outcomes summary
    outcomes_slide(p)

    # 11. Demo
    s = p.slides.add_slide(p.slide_layouts[6]); _bg(s, ACCENT)
    tf = textbox(s, 0.9, 2.2, 11.5, 3.2, anchor=MSO_ANCHOR.TOP)
    para(tf, "Live Demo", 40, WHITE, bold=True, first=True, after=18)
    para(tf, "Enter a patient → instant risk, a follow-up flag, and the factors behind it",
         20, CHIP, bullet=True, after=10)
    para(tf, "Governance view: model card, evaluation, fairness, explanations",
         20, CHIP, bullet=True, after=10)
    para(tf, "Monitoring: live dashboards and the drift report", 20, CHIP, bullet=True)

    out = HERE / "readmission_demo.pptx"
    p.save(out)
    print(f"wrote {out}  ({len(p.slides)} slides)")


def _down(slide, cx, top, height):
    a = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(cx - 0.11), Inches(top),
                               Inches(0.22), Inches(height))
    a.shadow.inherit = False; a.fill.solid(); a.fill.fore_color.rgb = DEEP
    a.line.fill.background()


def _eyebrow(slide, x, y, text):
    tf = textbox(slide, x, y, 8.0, 0.32)
    para(tf, text, 13, ACCENT, bold=True, first=True)


def architecture_slide(p):
    """Faithful, landscape recreation of the ARCHITECTURE.md flow (native shapes)."""
    s = p.slides.add_slide(p.slide_layouts[6]); _bg(s, BG)
    head(s, "Architecture & Data Flow")
    CX = 13.333 / 2

    # --- Offline plane: pipeline of boxes, each labelled with its output ---
    _eyebrow(s, 0.58, 1.05, "OFFLINE  —  build pipeline  (dvc repro → MLflow)")
    boxes = [
        ("Raw data", "diabetic_data.csv"),
        ("Clean", "cleaned.parquet"),
        ("Features", "features.parquet · 48"),
        ("Train + select", "XGB · LGBM · LR"),
        ("Registry", "model vN"),
    ]
    bw, bh, gap = 2.18, 0.95, 0.32
    total = len(boxes) * bw + (len(boxes) - 1) * gap
    x0 = (13.333 - total) / 2
    y = 1.45
    ends = []
    for i, (title, sub) in enumerate(boxes):
        x = x0 + i * (bw + gap)
        b = rect(s, x, y, bw, bh, ACCENT)
        tfx = b.text_frame; tfx.vertical_anchor = MSO_ANCHOR.MIDDLE
        tfx.margin_top = 0; tfx.margin_bottom = 0
        para(tfx, title, 13, WHITE, bold=True, first=True, align=PP_ALIGN.CENTER, after=1)
        para(tfx, sub, 9.5, CHIP, align=PP_ALIGN.CENTER, after=0)
        ends.append(x + bw)
    for i in range(len(boxes) - 1):
        a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(ends[i] + 0.02),
                               Inches(y + bh / 2 - 0.10), Inches(gap - 0.04), Inches(0.20))
        a.shadow.inherit = False; a.fill.solid(); a.fill.fore_color.rgb = DEEP
        a.line.fill.background()

    # --- bridge: export -> artifacts pill ---
    _down(s, CX, y + bh + 0.05, 0.40)
    elbl = textbox(s, CX + 0.18, y + bh + 0.06, 1.6, 0.3)
    para(elbl, "export_model", 10, MUTED, first=True)
    pill = rect(s, (13.333 - 7.2) / 2, 2.92, 7.2, 0.72, CHIP, line=LINE)
    tfp = pill.text_frame; tfp.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(tfp, "artifacts/   model.joblib  ·  feature_spec.json  ·  drift_reference.parquet",
         12.5, DEEP, bold=True, first=True, align=PP_ALIGN.CENTER)

    # --- Online plane: serve ---
    _down(s, CX, 3.70, 0.40)
    _eyebrow(s, 0.58, 4.18, "ONLINE  —  serve  (FastAPI in Docker)")
    serve = rect(s, (13.333 - 4.8) / 2, 4.20, 4.8, 0.9, ACCENT)
    tfs = serve.text_frame; tfs.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(tfs, "Serve / predict", 14, WHITE, bold=True, first=True, align=PP_ALIGN.CENTER, after=1)
    para(tfs, "featurize → model → risk + SHAP", 10, CHIP, align=PP_ALIGN.CENTER)

    # --- cross-cutting band ---
    _down(s, CX, 5.15, 0.38)
    band = rect(s, 0.58, 5.55, 12.17, 1.35, PANEL, line=LINE)
    tfb = band.text_frame; tfb.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(tfb, "Monitoring & Governance   (wraps the served model)", 15, ACCENT,
         bold=True, first=True, align=PP_ALIGN.CENTER, after=4)
    para(tfb, "every prediction logged  →  drift · metrics · alerts      |      "
         "fairness · SHAP · model card · lineage", 12.5, TEXT, align=PP_ALIGN.CENTER)
    return s


def modeling_slide(p):
    s = p.slides.add_slide(p.slide_layouts[6]); _bg(s, BG)
    head(s, "Stage 3 — Modeling & Evaluation")
    tf = textbox(s, 0.6, 1.15, 12.1, 1.75)
    para(tf, "DECISIONS", 13, ACCENT, bold=True, first=True, after=6)
    para(tf, "Split by patient (no leakage) · tune on PR-AUC · calibrate probabilities · "
         "cost-based threshold (0.10).", 14, TEXT, after=5)
    para(tf, "Handle 11% imbalance by weighting the minority class ≈7.8× in the loss "
         "(scale_pos_weight = negatives ÷ positives = 61,698 ÷ 7,915) — not resampling, "
         "which would duplicate patients and break grouped CV.", 13, MUTED)

    rows = [
        ["Model (validation)", "PR-AUC", "ROC-AUC", "Recall@p≥30%", "Brier"],
        ["Logistic regression (baseline)", "0.221", "0.667", "0.189", "0.225"],
        ["XGBoost  — winner", "0.239", "0.680", "0.216", "0.206"],
        ["LightGBM", "0.239", "0.675", "0.215", "0.201"],
        ["XGBoost + calibration", "0.236", "0.680", "0.205", "0.097"],
    ]
    col_w = [Inches(4.0), Inches(1.5), Inches(1.5), Inches(2.0), Inches(1.5)]
    gt = s.shapes.add_table(len(rows), len(rows[0]), Inches(1.0), Inches(2.95),
                            Inches(10.5), Inches(0.45 * len(rows))).table
    for j, cw in enumerate(col_w):
        gt.columns[j].width = cw
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            c = gt.cell(i, j); c.text = str(val)
            pa = c.text_frame.paragraphs[0]
            pa.font.size = Pt(13); pa.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            c.fill.solid()
            if i == 0:
                c.fill.fore_color.rgb = ACCENT; pa.font.color.rgb = WHITE; pa.font.bold = True
            else:
                hi = (i == 2)
                c.fill.fore_color.rgb = CHIP if hi else WHITE
                pa.font.color.rgb = TEXT; pa.font.bold = hi

    tf2 = textbox(s, 0.6, 5.7, 12.1, 1.3)
    para(tf2, f"OUTCOME — one-time test: recall {TA['recall']:.0%} of readmissions caught, "
         f"PR-AUC {T['test_pr_auc']:.3f}, Brier {T['test_brier']:.3f} (calibration fixed "
         "0.206 → 0.097).", 15, TEXT, bold=True, first=True, after=6)
    para(tf2, f"Threshold 0.10 from real costs: flag when risk × {CM['preventable_fraction']} "
         f"× ${CM['readmission_cost']:,} > ${CM['intervention_cost']:,}. Low precision is "
         "deliberate — a cheap call vs a costly readmission.", 13, MUTED)
    return s


def outcomes_slide(p):
    s = p.slides.add_slide(p.slide_layouts[6]); _bg(s, BG)
    head(s, "Outcomes & Honest Limitations")
    # three metric tiles
    tiles = [
        (f"{TA['recall']:.0%}", "of 30-day readmissions flagged"),
        (f"${TA['net_savings_per_1000']:,.0f}", "net savings per 1,000 patients"),
        (f"{T['test_brier']:.3f}", "Brier — calibrated, honest %"),
    ]
    tw, gap = 3.7, 0.5
    x0 = (13.333 - (3 * tw + 2 * gap)) / 2
    for i, (big, lab) in enumerate(tiles):
        x = x0 + i * (tw + gap)
        rect(s, x, 1.35, tw, 1.7, WHITE, line=LINE)
        tf = textbox(s, x + 0.1, 1.45, tw - 0.2, 1.5, anchor=MSO_ANCHOR.MIDDLE)
        para(tf, big, 34, ACCENT, bold=True, first=True, align=PP_ALIGN.CENTER, after=2)
        para(tf, lab, 13, MUTED, align=PP_ALIGN.CENTER)

    tf = textbox(s, 0.6, 3.5, 6.0, 3.4)
    para(tf, "WHAT WORKS", 14, ACCENT, bold=True, first=True, after=8)
    for it in ["End-to-end, reproducible pipeline (data → governance)",
               "Calibrated, explainable, cost-aligned predictions",
               "Live monitoring, drift, and a retraining trigger",
               "Fairness audited and disclosed; model card from code"]:
        para(tf, it, 14, TEXT, bullet=True, after=7)

    tf2 = textbox(s, 6.9, 3.5, 5.8, 3.4)
    para(tf2, "LIMITATIONS", 14, ACCENT, bold=True, first=True, after=8)
    for it in ["1999–2008 / ICD-9 data — retrain before real use",
               "Low precision by design (high-recall triage)",
               "Weaker for young & rare subgroups",
               "Predicts risk, not clinical need"]:
        para(tf2, it, 14, TEXT, bullet=True, after=7)
    return s


if __name__ == "__main__":
    build()
